import streamlit as st
import fitz  # PyMuPDF
from pptx import Presentation
from docx import Document
import difflib
import matplotlib.pyplot as plt
import numpy as np
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
from wordcloud import WordCloud
from collections import Counter
import seaborn as sns
import pandas as pd
import io
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from reportlab.lib.utils import ImageReader
import base64  # Import base64 module for encoding binary data


# Function to read text from different file types
def read_file(file):
    content = ""
    try:
        file_extension = file.name.split('.')[-1].lower()
        if file_extension == 'pdf':
            pdf_document = fitz.open(stream=file.read(), filetype="pdf")
            for page_num in range(len(pdf_document)):
                page = pdf_document.load_page(page_num)
                content += page.get_text() + "\n"
        elif file_extension == 'pptx':
            ppt = Presentation(io.BytesIO(file.read()))
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        content += shape.text + "\n"
        elif file_extension == 'docx':
            doc = Document(io.BytesIO(file.read()))
            for para in doc.paragraphs:
                content += para.text + "\n"
        else:
            content = file.read().decode('utf-8')  # Assume UTF-8 encoding for other file types
    except Exception as e:
        st.error(f"Error reading file: {e}")
        content = None
    return content


# Function to compare texts and return similarity percentage and differences
def compare_text(text1, text2):
    try:
        d = difflib.SequenceMatcher(None, text1, text2)
        similarity_ratio = d.ratio()
        similarity_percentage = int(similarity_ratio * 100)
        diff = list(d.get_opcodes())
        return similarity_percentage, diff
    except Exception as e:
        st.error(f"Error comparing texts: {e}")
        return None, None


# Function to clean and tokenize text
def clean_and_tokenize(text):
    # Tokenization and basic cleaning (lowercase, remove punctuation, numbers, etc.)
    tokens = text.lower().split()
    tokens = [token.strip(",.!?") for token in tokens if token.isalpha()]
    return tokens


# Function to plot word frequency and return plot image
def plot_word_frequency(text1, text2):
    try:
        # Clean and tokenize text, removing common words
        common_words = set(['and', 'or', 'not', 'the', 'is', 'in', 'of', 'to', 'a', 'on'])
        tokens1 = clean_and_tokenize(text1)
        tokens2 = clean_and_tokenize(text2)

        tokens1 = [token for token in tokens1 if token not in common_words]
        tokens2 = [token for token in tokens2 if token not in common_words]

        # Calculate word frequencies
        word_freq1 = Counter(tokens1)
        word_freq2 = Counter(tokens2)

        # Get top 20 frequent words
        top_words1, top_freqs1 = zip(*word_freq1.most_common(20))
        top_words2, top_freqs2 = zip(*word_freq2.most_common(20))

        # Plotting
        fig, axes = plt.subplots(2, 1, figsize=(10, 12))

        # Plot for Text 1
        sns.barplot(ax=axes[0], x=top_freqs1, y=top_words1, palette="Blues_d")
        axes[0].set_title('Top 20 Words - Text 1')
        axes[0].set_xlabel('Frequency')
        axes[0].set_ylabel('Words')
        axes[0].set_xlim(0, max(top_freqs1) + 1)
        axes[0].text(-0.1, 1.1, "A bar plot showing the top 20 most frequent words in Text 1. Each bar represents a word and its height represents its frequency in the text.", fontsize=10, transform=axes[0].transAxes, verticalalignment='top', bbox=dict(facecolor='white', alpha=0.5))

        # Plot for Text 2
        sns.barplot(ax=axes[1], x=top_freqs2, y=top_words2, palette="Greens_d")
        axes[1].set_title('Top 20 Words - Text 2')
        axes[1].set_xlabel('Frequency')
        axes[1].set_ylabel('Words')
        axes[1].set_xlim(0, max(top_freqs2) + 1)
        axes[1].text(-0.1, 1.1, "A bar plot showing the top 20 most frequent words in Text 2. Each bar represents a word and its height represents its frequency in the text.", fontsize=10, transform=axes[1].transAxes, verticalalignment='top', bbox=dict(facecolor='white', alpha=0.5))

        # Adjust subplot layout
        plt.tight_layout()

        # Convert plot to PNG image
        word_freq_bytes = io.BytesIO()
        plt.savefig(word_freq_bytes, format='png')
        plt.close()

        return word_freq_bytes

    except Exception as e:
        st.error(f"Error plotting word frequency: {e}")
        return None


# Function to plot cosine similarity heatmap and return plot image
def plot_cosine_similarity(text1, text2):
    try:
        # Vectorize texts using TF-IDF
        vectorizer = TfidfVectorizer()
        X = vectorizer.fit_transform([text1, text2])

        # Calculate cosine similarity
        cosine_sim = cosine_similarity(X)

        # Plotting heatmap
        fig, ax = plt.subplots(figsize=(8, 6))
        sns.heatmap(cosine_sim, annot=True, cmap="YlGnBu", xticklabels=['Text 1', 'Text 2'], yticklabels=['Text 1', 'Text 2'], ax=ax)
        ax.set_title('Cosine Similarity Heatmap')
        ax.set_xlabel('Text')
        ax.set_ylabel('Text')
        ax.text(-0.1, 1.1, "A heatmap showing the cosine similarity between Text 1 and Text 2. Higher values indicate greater similarity, while lower values indicate less similarity between the texts.", fontsize=10, transform=ax.transAxes, verticalalignment='top', bbox=dict(facecolor='white', alpha=0.5))

        # Convert plot to PNG image
        cosine_sim_bytes = io.BytesIO()
        plt.savefig(cosine_sim_bytes, format='png')
        plt.close()

        return cosine_sim_bytes

    except Exception as e:
        st.error(f"Error plotting cosine similarity: {e}")
        return None


# Function to generate word cloud and return plot image
def generate_word_cloud(text):
    try:
        # Generate word cloud
        wordcloud = WordCloud(width=800, height=400, background_color='white').generate(text)

        # Plotting word cloud
        fig, ax = plt.subplots(figsize=(10, 6))
        ax.imshow(wordcloud, interpolation='bilinear')
        ax.axis('off')
        ax.set_title('Word Cloud')
        ax.text(-0.1, 1.1, "A word cloud visualization representing the frequency of words in the text. Larger words indicate higher frequency.", fontsize=10, transform=ax.transAxes, verticalalignment='top', bbox=dict(facecolor='white', alpha=0.5))

        # Convert plot to PNG image
        word_cloud_bytes = io.BytesIO()
        plt.savefig(word_cloud_bytes, format='png')
        plt.close()

        return word_cloud_bytes

    except Exception as e:
        st.error(f"Error generating word cloud: {e}")
        return None


# Function to plot scatter plot of word frequencies and return plot image
def plot_scatter(text1, text2):
    try:
        # Clean and tokenize text, removing common words
        common_words = set(['and', 'or', 'not', 'the', 'is', 'in', 'of', 'to', 'a', 'on'])
        tokens1 = clean_and_tokenize(text1)
        tokens2 = clean_and_tokenize(text2)

        tokens1 = [token for token in tokens1 if token not in common_words]
        tokens2 = [token for token in tokens2 if token not in common_words]

        # Calculate word frequencies
        word_freq1 = Counter(tokens1)
        word_freq2 = Counter(tokens2)

        # Create DataFrame
        freq_df = pd.DataFrame(list(set(word_freq1.keys()).union(set(word_freq2.keys()))), columns=['word'])
        freq_df['freq_text1'] = freq_df['word'].apply(lambda x: word_freq1.get(x, 0))
        freq_df['freq_text2'] = freq_df['word'].apply(lambda x: word_freq2.get(x, 0))

        # Plotting scatter plot
        fig, ax = plt.subplots(figsize=(10, 6))
        ax.scatter(freq_df['freq_text1'], freq_df['freq_text2'], alpha=0.5)
        ax.set_xlabel('Word Frequency in Text 1')
        ax.set_ylabel('Word Frequency in Text 2')
        ax.set_title('Scatter Plot of Word Frequencies')
        ax.grid(True)
        ax.text(-0.1, 1.1, "A scatter plot comparing the word frequencies between Text 1 and Text 2. Each point represents a word and its position indicates its frequency in both texts.", fontsize=10, transform=ax.transAxes, verticalalignment='top', bbox=dict(facecolor='white', alpha=0.5))

        # Convert plot to PNG image
        scatter_plot_bytes = io.BytesIO()
        plt.savefig(scatter_plot_bytes, format='png')
        plt.close()

        return scatter_plot_bytes

    except Exception as e:
        st.error(f"Error plotting scatter plot: {e}")
        return None


# Streamlit app
def main():
    st.title("Text Comparison and Plagiarism Checker")
    
    # File upload and text display
    st.sidebar.title("Upload Files")
    uploaded_file1 = st.sidebar.file_uploader("Upload File 1", type=['pdf', 'pptx', 'docx', 'txt'])
    uploaded_file2 = st.sidebar.file_uploader("Upload File 2", type=['pdf', 'pptx', 'docx', 'txt'])

    if uploaded_file1 and uploaded_file2:
        st.sidebar.subheader("File 1")
        file_details1 = {"FileName": uploaded_file1.name, "FileType": uploaded_file1.type, "FileSize": uploaded_file1.size}
        st.sidebar.write(file_details1)

        st.sidebar.subheader("File 2")
        file_details2 = {"FileName": uploaded_file2.name, "FileType": uploaded_file2.type, "FileSize": uploaded_file2.size}
        st.sidebar.write(file_details2)

        # Read files
        text1 = read_file(uploaded_file1)
        text2 = read_file(uploaded_file2)

        if text1 and text2:
            # Compare texts
            similarity_percentage, diff = compare_text(text1, text2)

            if similarity_percentage is not None:
                st.subheader("Similarity Analysis")
                st.write(f"Similarity Percentage: {similarity_percentage}%")

                # Generate plots and PDF report
                st.subheader("Plots and PDF Report")
                st.write("Below are the visualizations and PDF report:")
                
                # Generate word frequency plot
                word_freq_image = plot_word_frequency(text1, text2)
                if word_freq_image:
                    st.image(word_freq_image, use_column_width=True, caption='Word Frequency Plot')

                # Generate cosine similarity heatmap
                cosine_sim_image = plot_cosine_similarity(text1, text2)
                if cosine_sim_image:
                    st.image(cosine_sim_image, use_column_width=True, caption='Cosine Similarity Heatmap')

                # Generate word cloud for Text 1
                word_cloud_image1 = generate_word_cloud(text1)
                if word_cloud_image1:
                    st.image(word_cloud_image1, use_column_width=True, caption='Word Cloud - Text 1')

                # Generate word cloud for Text 2
                word_cloud_image2 = generate_word_cloud(text2)
                if word_cloud_image2:
                    st.image(word_cloud_image2, use_column_width=True, caption='Word Cloud - Text 2')

                # Generate scatter plot of word frequencies
                scatter_plot_image = plot_scatter(text1, text2)
                if scatter_plot_image:
                    st.image(scatter_plot_image, use_column_width=True, caption='Scatter Plot of Word Frequencies')

                # Generate PDF report
                st.markdown("*")
                st.subheader("Download PDF Report")
                pdf_report = generate_pdf_report(text1, text2, similarity_percentage)
                if pdf_report:
                    st.markdown(get_binary_file_downloader_html(pdf_report, "Download PDF", "Report"), unsafe_allow_html=True)

            else:
                st.warning("Error occurred during text comparison.")
        else:
            st.warning("Error reading one or both files.")
    else:
        st.info("Please upload two files to perform comparison.")


# Function to generate PDF report
def generate_pdf_report(text1, text2, similarity_percentage):
    try:
        # Create a canvas
        pdf_buffer = io.BytesIO()
        c = canvas.Canvas(pdf_buffer, pagesize=letter)
        
        # Set up title
        c.setFont("Helvetica-Bold", 16)
        c.drawString(100, 750, "Text Comparison Report")

        # Display similarity percentage
        c.setFont("Helvetica", 12)
        c.drawString(100, 730, f"Similarity Percentage: {similarity_percentage}%")

        # Generate word cloud and add to PDF
        word_cloud_image1 = generate_word_cloud(text1)
        if word_cloud_image1:
            word_cloud_image1.seek(0)
            c.drawImage(ImageReader(word_cloud_image1), 100, 500, width=400, height=300, mask='auto')

        word_cloud_image2 = generate_word_cloud(text2)
        if word_cloud_image2:
            word_cloud_image2.seek(0)
            c.drawImage(ImageReader(word_cloud_image2), 100, 150, width=400, height=300, mask='auto')

        # Save the PDF
        c.showPage()  # Start a new page for additional plots

        # Generate word frequency plot and add to PDF
        word_freq_image = plot_word_frequency(text1, text2)
        if word_freq_image:
            word_freq_image.seek(0)
            c.drawImage(ImageReader(word_freq_image), 100, 500, width=400, height=300, mask='auto')

        # Generate cosine similarity heatmap and add to PDF
        cosine_sim_image = plot_cosine_similarity(text1, text2)
        if cosine_sim_image:
            cosine_sim_image.seek(0)
            c.drawImage(ImageReader(cosine_sim_image), 100, 150, width=400, height=300, mask='auto')

        # Save the PDF
        c.save()
        pdf_buffer.seek(0)
        
        return pdf_buffer

    except Exception as e:
        st.error(f"Error generating PDF report: {e}")
        return None


# Function to provide a download link for generated files
def get_binary_file_downloader_html(bin_file, file_label='File', button_label='Download'):
    b64 = base64.b64encode(bin_file.getvalue()).decode()  # Encode the PDF buffer in base64
    href = f'<a href="data:application/octet-stream;base64,{b64}" download="{file_label}.pdf">{button_label}</a>'
    return href


if __name__ == "__main__":
    main()